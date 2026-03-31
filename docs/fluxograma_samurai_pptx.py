#!/usr/bin/env python3
"""Gera fluxograma PPTX do processo SamurAI - BlackBelt Platform"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BLUE = RGBColor(0x1e, 0x40, 0xaf)
GREEN = RGBColor(0x16, 0xa3, 0x4a)
ORANGE = RGBColor(0xea, 0x58, 0x0c)
RED = RGBColor(0xdc, 0x26, 0x26)
PURPLE = RGBColor(0x7c, 0x3a, 0xed)
GOLD = RGBColor(0xd9, 0x77, 0x06)
DARK = RGBColor(0x1f, 0x29, 0x37)
GRAY = RGBColor(0x6b, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BLUE_LIGHT = RGBColor(0xdb, 0xea, 0xfe)
GREEN_LIGHT = RGBColor(0xdc, 0xfc, 0xe7)
ORANGE_LIGHT = RGBColor(0xff, 0xf7, 0xed)
RED_LIGHT = RGBColor(0xfe, 0xf2, 0xf2)
PURPLE_LIGHT = RGBColor(0xf5, 0xf3, 0xff)
GOLD_LIGHT = RGBColor(0xff, 0xfb, 0xeb)

def add_step(slide, left, top, width, height, number, title, subtitle, bg_color, border_color):
    """Add a step box to the slide"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Number + Title
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{number}. {title}"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = DARK
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(7)
    run2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.LEFT

def add_arrow_right(slide, x, y, length):
    """Add right arrow"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, length, Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()

def add_arrow_down(slide, x, y, length):
    """Add down arrow"""
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.2), length)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # === SLIDE 1: Title ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xf8, 0xfa, 0xfc)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.3), Inches(12.73), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.line.fill.background()

    tf = header.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "SamurAI - Fluxo Completo NR-01"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Black Belt Consultoria | Gestao de Riscos Psicossociais"
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(0xbe, 0xd0, 0xfe)
    p2.alignment = PP_ALIGN.CENTER

    # Steps layout
    sw = Inches(2.2)
    sh = Inches(0.85)
    gap_x = Inches(0.25)
    gap_y = Inches(0.15)

    # ROW 1 (y=1.8)
    row1_y = Inches(1.8)
    steps_r1 = [
        (1, "Inserir CNPJ", "Consultor digita CNPJ\n+ funcionarios + email", BLUE_LIGHT, BLUE),
        (2, "Consulta Receita", "Busca automatica\nna Receita Federal", BLUE_LIGHT, BLUE),
        (3, "Cadastro Empresa", "Cria tenant + checklist\nNR-01 automatico", GREEN_LIGHT, GREEN),
        (4, "Pre-Proposta", "Gera estimativa com\nbase no porte/setor", ORANGE_LIGHT, ORANGE),
        (5, "Editar Proposta", "Modal visual para\najustar valores/itens", ORANGE_LIGHT, ORANGE),
    ]

    for i, (n, t, s, bg, bc) in enumerate(steps_r1):
        x = Inches(0.3) + i * (sw + gap_x)
        add_step(slide, x, row1_y, sw, sh, n, t, s, bg, bc)
        if i < 4:
            add_arrow_right(slide, x + sw, row1_y + sh/2 - Inches(0.1), gap_x)

    # Arrow down
    add_arrow_down(slide, Inches(0.3) + 4*(sw+gap_x) + sw/2 - Inches(0.1), row1_y + sh, gap_y)

    # ROW 2 (y=2.8) - right to left
    row2_y = Inches(2.8)
    steps_r2 = [
        (10, "Enviar COPSOQ", "Dispara questionarios\npor email", RED_LIGHT, RED),
        (9, "Cadastrar Pessoas", "Empresa cadastra setores\ne colaboradores", PURPLE_LIGHT, PURPLE),
        (8, "Criar Conta", "Auto-cria usuario +\nenvia credenciais", PURPLE_LIGHT, PURPLE),
        (7, "Aprovacao Email", "Empresa clica Aprovar\nno email recebido", GREEN_LIGHT, GREEN),
        (6, "Enviar por Email", "Email com proposta\n+ botoes Aprovar/Recusar", GREEN_LIGHT, GREEN),
    ]

    for i, (n, t, s, bg, bc) in enumerate(steps_r2):
        x = Inches(0.3) + i * (sw + gap_x)
        add_step(slide, x, row2_y, sw, sh, n, t, s, bg, bc)
        if i < 4:
            add_arrow_right(slide, x + sw, row2_y + sh/2 - Inches(0.1), gap_x)

    # Arrow down
    add_arrow_down(slide, Inches(0.3) + sw/2 - Inches(0.1), row2_y + sh, gap_y)

    # ROW 3 (y=3.8)
    row3_y = Inches(3.8)
    steps_r3 = [
        (11, "Respostas COPSOQ", "Colaboradores respondem\n(anonimo, 7 dias)", RED_LIGHT, RED),
        (12, "Gerar Relatorio", "Analise automatica\ndos resultados", GOLD_LIGHT, GOLD),
        (13, "Inventario Riscos", "Mapeamento completo\nde riscos psicossociais", GOLD_LIGHT, GOLD),
        (14, "Plano de Acao", "Acoes corretivas com\nprazos e responsaveis", GOLD_LIGHT, GOLD),
        (15, "Proposta Final", "Proposta definitiva com\nvalores reais do servico", GREEN_LIGHT, GREEN),
    ]

    for i, (n, t, s, bg, bc) in enumerate(steps_r3):
        x = Inches(0.3) + i * (sw + gap_x)
        add_step(slide, x, row3_y, sw, sh, n, t, s, bg, bc)
        if i < 4:
            add_arrow_right(slide, x + sw, row3_y + sh/2 - Inches(0.1), gap_x)

    # Arrow down
    add_arrow_down(slide, Inches(0.3) + 4*(sw+gap_x) + sw/2 - Inches(0.1), row3_y + sh, gap_y)

    # ROW 4 (y=4.8)
    row4_y = Inches(4.8)
    steps_r4 = [
        (18, "Entrega Final", "7 PDFs entregues:\nProposta, COPSOQ, etc.", BLUE_LIGHT, BLUE),
        (17, "Liberar PDFs", "Todos os documentos\nliberados para download", BLUE_LIGHT, BLUE),
        (16, "Aprovar + Pagar", "Empresa aprova proposta\nfinal e paga parcelas", GREEN_LIGHT, GREEN),
    ]

    for i, (n, t, s, bg, bc) in enumerate(steps_r4):
        x = Inches(0.3) + (4-i) * (sw + gap_x)
        add_step(slide, x, row4_y, sw, sh, n, t, s, bg, bc)
        if i < 2:
            prev_x = Inches(0.3) + (4-i) * (sw + gap_x)
            next_x = Inches(0.3) + (4-i-1) * (sw + gap_x)

    # Deliverables box
    dx = Inches(0.3)
    dy = Inches(5.85)
    dbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dx, dy, Inches(12.73), Inches(0.8))
    dbox.fill.solid()
    dbox.fill.fore_color.rgb = RGBColor(0xf3, 0xf4, 0xf6)
    dbox.line.color.rgb = GRAY
    dbox.line.width = Pt(1)

    tf = dbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Documentos Entregues (7 PDFs)"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Proposta Comercial | Relatorio COPSOQ-II | Inventario de Riscos | Plano de Acao | Programa de Treinamento | Checklist de Conformidade | Certificado NR-01"
    run2.font.size = Pt(9)
    run2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER

    # Legend
    legend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(6.8), Inches(12.73), Inches(0.45))
    legend.fill.solid()
    legend.fill.fore_color.rgb = WHITE
    legend.line.fill.background()

    tf = legend.text_frame
    p = tf.paragraphs[0]
    items = [
        ("Consultoria", BLUE), ("Aprovacao/Envio", GREEN),
        ("Proposta", ORANGE), ("Empresa", PURPLE),
        ("COPSOQ-II", RED), ("Analise NR-01", GOLD),
    ]
    for i, (label, color) in enumerate(items):
        run = p.add_run()
        run.text = f"  ● {label}  "
        run.font.size = Pt(9)
        run.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(12.73), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Black Belt Consultoria SST | contato@blackbeltconsultoria.com | blackbeltconsultoria.com"
    run.font.size = Pt(8)
    run.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    prs.save("docs/Fluxograma_SamurAI_BlackBelt.pptx")
    print("PPTX criado: docs/Fluxograma_SamurAI_BlackBelt.pptx")

if __name__ == "__main__":
    create_pptx()
